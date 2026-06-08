import os, re, json, csv, io
from typing import List, Tuple, Dict, Optional, Union
from collections import Counter
from pathlib import Path

from backend.token_counter import count_tokens
from backend.prompt_optimizer import OptiTokenOptimizer

# Precompile register detection patterns for performance
_REGISTER_KW_PATTERNS: Dict[str, re.Pattern] = {}
_REGISTER_SENTINEL_PATTERNS: Dict[str, re.Pattern] = {}

REGISTER_KEYWORDS = {
    "scientific": [
        "study", "research", "experiment", "analysis", "data", "hypothesis",
        "methodology", "results", "findings", "conclusion", "correlation",
        "significant", "sample", "variable", "measurement", "calibration",
        "protocol", "laboratory", "specimen", "quantitative", "qualitative",
        "étude", "recherche", "expérience", "analyse", "données", "hypothèse",
        "méthodologie", "résultats", "conclusion", "corrélation",
    ],
    "legal": [
        "contract", "agreement", "clause", "party", "obligation", "liability",
        "indemnify", "termination", "breach", "warrant", "governing law",
        "jurisdiction", "arbitration", "confidentiality", "hereby",
        "notwithstanding", "pursuant", "thereof", "therein", "wherein",
        "contrat", "convention", "clause", "partie", "obligation", "responsabilité",
        "garantir", "résiliation", "manquement", "droit applicable",
    ],
    "financial": [
        "revenue", "profit", "loss", "asset", "liability", "equity",
        "balance sheet", "income statement", "cash flow", "depreciation",
        "amortization", "dividend", "shareholder", "fiscal", "audit",
        "budget", "forecast", "margin", "earnings", "ebitda",
        "chiffre d'affaires", "bénéfice", "perte", "actif", "passif",
        "bilan", "compte de résultat", "trésorerie", "dividende",
    ],
    "technical": [
        "specification", "requirement", "architecture", "interface", "protocol",
        "implementation", "deployment", "configuration", "parameter",
        "threshold", "throughput", "latency", "bandwidth", "redundancy",
        "spécification", "exigence", "architecture", "interface", "protocole",
        "implémentation", "déploiement", "configuration", "paramètre",
    ],
    "administrative": [
        "policy", "procedure", "regulation", "compliance", "guideline",
        "directive", "memorandum", "circular", "form", "submission",
        "approval", "authorization", "certification", "standard",
        "politique", "procédure", "règlement", "conformité", "directive",
        "note de service", "circulaire", "soumission", "approbation",
    ],
    "commercial": [
        "marketing", "campaign", "product", "customer", "client", "brand",
        "market", "strategy", "promotion", "offer", "discount", "loyalty",
        "acquisition", "conversion", "pipeline", "prospect", "engagement",
        "marché", "campagne", "produit", "client", "marque", "stratégie",
        "promotion", "offre", "remise", "fidélité",
    ],
    "academic": [
        "paper", "publication", "citation", "reference", "bibliography",
        "abstract", "introduction", "literature review", "appendix",
        "thesis", "dissertation", "peer review", "journal", "conference",
        "article", "publication", "citation", "référence", "bibliographie",
        "résumé", "introduction", "revue de littérature", "annexe",
        "thèse", "mémoire", "comité de lecture",
    ],
    "literary": [
        "chapter", "scene", "narrative", "protagonist", "character",
        "dialogue", "metaphor", "poem", "verse", "prose", "fiction",
        "chapitre", "scène", "récit", "protagoniste", "personnage",
        "dialogue", "métaphore", "poème", "vers", "prose", "fiction",
    ],
    "philosophical": [
        "therefore", "hence", "thus", "premise", "conclusion", "argument",
        "logic", "reason", "essence", "existence", "cause", "effect",
        "truth", "knowledge", "consciousness", "paradox", "dialectic",
        "donc", "car", "or", "prémisse", "conclusion", "argument",
        "logique", "raison", "essence", "existence", "cause", "effet",
        "vérité", "connaissance", "conscience", "paradoxe", "dialectique",
    ],
    "instructional": [
        "step", "steps", "tutorial", "guide", "instructions", "procedure",
        "how to", "follow", "first", "second", "next", "then", "finally",
        "étape", "étapes", "tutoriel", "guide", "instructions", "procédure",
        "comment", "suivez", "premièrement", "deuxièmement", "enfin",
    ],
}

REGISTER_SENTINELS = {
    "scientific": ["fig.", "table", "eq.", "p<0.05", "p<0.01", "n=", "±", "σ", "α="],
    "legal": ["§", "art.", "article", "section", "subsection", "c.", "v."],
    "financial": ["$", "€", "£", "%", "k$", "m$", "b$", "taux", "ratio"],
    "technical": ["api", "json", "xml", "tcp", "http", "https", "ftp", "ssh"],
    "philosophical": ["→", "∴", "∵", "∀", "∃", "¬", "⇒", "⇔"],
    "instructional": ["1.", "2.", "3.", "step 1", "étape 1", "1ère étape"],
}

# Build precompiled patterns once at module level
for _reg, _kws in REGISTER_KEYWORDS.items():
    _escaped = [re.escape(kw) for kw in _kws]
    _REGISTER_KW_PATTERNS[_reg] = re.compile('|'.join(_escaped), re.I)
for _reg, _sents in REGISTER_SENTINELS.items():
    _escaped = [re.escape(s) for s in _sents]
    _REGISTER_SENTINEL_PATTERNS[_reg] = re.compile('|'.join(_escaped))

class DocumentContent:
    """Représentation structurée du contenu extrait d'un document."""
    def __init__(self):
        self.title: str = ""
        self.metadata: Dict[str, str] = {}
        self.sections: List[Dict] = []       # [{heading, level, paragraphs: [], tables: [], lists: []}]
        self.paragraphs: List[str] = []
        self.tables: List[Dict] = []          # [{headers: [], rows: [[]]}]
        self.lists: List[List[str]] = []
        self.raw_text: str = ""
        self.page_count: int = 0
        self.detected_format: str = ""
        self.filename: str = ""

    def to_text(self, include_tables: bool = True, include_lists: bool = True) -> str:
        """Reconstruit le texte structuré."""
        parts = []
        if self.title:
            parts.append(f"# {self.title}")
        for sec in self.sections:
            if sec.get("heading"):
                prefix = "#" * min(sec.get("level", 1), 6)
                parts.append(f"\n{prefix} {sec['heading']}")
            for p in sec.get("paragraphs", []):
                parts.append(p)
            if include_tables:
                for t in sec.get("tables", []):
                    parts.append(self._table_to_text(t))
            if include_lists:
                for lst in sec.get("lists", []):
                    for i, item in enumerate(lst, 1):
                        parts.append(f"  {i}. {item}")
        # Fallback: paragraphs not in sections
        if not self.sections:
            for p in self.paragraphs:
                parts.append(p)
            if include_tables:
                for t in self.tables:
                    parts.append(self._table_to_text(t))
            if include_lists:
                for lst in self.lists:
                    for i, item in enumerate(lst, 1):
                        parts.append(f"  {i}. {item}")
        return "\n".join(parts)

    @staticmethod
    def _table_to_text(t: Dict) -> str:
        headers = t.get("headers", [])
        rows = t.get("rows", [])
        if not headers and not rows:
            return ""
        lines = []
        if headers:
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join("---" for _ in headers) + " |")
        for row in rows:
            lines.append("| " + " | ".join(str(c) for c in row) + " |")
        return "\n".join(lines)


class DocumentAnalyzer:
    """
    Analyse et comprime des documents (docx, pdf, pptx, xlsx, txt, md, odt, html, csv, json, xml).
    Pipeline: Parse → Extraction structure → Détection registre → Compression (Light/Aggressive)
    """

    # Parser dispatch
    TEXT_EXTENSIONS = {".txt", ".md", ".rst", ".log", ".ini", ".cfg", ".conf", ".yml", ".yaml", ".toml"}
    OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".odt", ".ods", ".odp", ".rtf"}
    WEB_EXTENSIONS = {".html", ".htm", ".xhtml"}
    DATA_EXTENSIONS = {".csv", ".json", ".xml", ".tsv"}

    def __init__(self, model: str = "gpt-4o"):
        self.model = model
        self.optimizer = OptiTokenOptimizer()

    # ──────────────────────────────────────────────
    #  PUBLIC API
    # ──────────────────────────────────────────────

    def analyze(self, filepath: Union[str, Path], filename: str = "") -> DocumentContent:
        """Parse un fichier et retourne son contenu structuré."""
        path = Path(filepath)
        fname = filename or path.name
        ext = Path(fname).suffix.lower()
        content = DocumentContent()
        content.filename = fname
        content.detected_format = ext.lstrip(".")
        parser = self._get_parser(ext)
        if parser:
            parser(path, content)
        else:
            content.raw_text = path.read_text(encoding="utf-8", errors="replace")
            content.paragraphs = [p for p in content.raw_text.split("\n") if p.strip()]
        return content

    def analyze_bytes(self, data: bytes, filename: str) -> DocumentContent:
        """Parse depuis un buffer mémoire."""
        ext = Path(filename).suffix.lower()
        content = DocumentContent()
        content.filename = filename
        content.detected_format = ext.lstrip(".")
        parser = self._get_parser(ext)
        if parser:
            from tempfile import NamedTemporaryFile
            with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(data)
                tmp_path = tmp.name
            try:
                parser(Path(tmp_path), content)
            finally:
                os.unlink(tmp_path)
        else:
            content.raw_text = data.decode("utf-8", errors="replace")
            content.paragraphs = [p for p in content.raw_text.split("\n") if p.strip()]
        return content

    def compress(self, content: DocumentContent, mode: str = "light",
                 category: str = None) -> Tuple[str, Dict]:
        """
        Compresse le contenu d'un document.
        mode: 'light' (structure + Optimizer Light), 'balanced' (structure + Optimizer Balanced),
              'aggressive' (key sentences + Optimizer Aggressive).
        Retourne (texte_compressé, metadata).
        """
        if mode not in ("light", "balanced", "aggressive", "max", "industrial"):
            raise ValueError(f"Mode must be 'light', 'balanced', 'aggressive', 'max', or 'industrial', got '{mode}'")

        if category is None:
            category = self._detect_register(content.raw_text)

        metadata = {
            "mode": mode,
            "category": category,
            "original_tokens": count_tokens(content.raw_text, self.model),
            "sections_before": len(content.sections) or 1,
            "paragraphs_before": len(content.paragraphs) or sum(len(s.get("paragraphs",[])) for s in content.sections) or 1,
        }

        if mode == "light":
            compressed_text = self._compress_light(content, category, mode="light")
        elif mode == "balanced":
            compressed_text = self._compress_light(content, category, mode="balanced")
        elif mode == "aggressive":
            compressed_text = self._compress_aggressive(content, category)
        else:
            compressed_text = self._compress_spc(content, category, mode=mode)

        metadata["compressed_tokens"] = count_tokens(compressed_text, self.model)
        metadata["tokens_saved"] = metadata["original_tokens"] - metadata["compressed_tokens"]
        metadata["savings_percent"] = round(
            (metadata["tokens_saved"] / max(metadata["original_tokens"], 1)) * 100, 1
        )
        return compressed_text, metadata

    # ──────────────────────────────────────────────
    #  FORMAT DETECTION
    # ──────────────────────────────────────────────

    def _get_parser(self, ext: str):
        ext = ext.lower()
        if ext in self.TEXT_EXTENSIONS:
            return self._parse_text
        if ext == ".docx":
            return self._parse_docx
        if ext in {".pptx", ".ppsx"}:
            return self._parse_pptx
        if ext in {".xlsx", ".xlsm", ".xltx"}:
            return self._parse_xlsx
        if ext == ".pdf":
            return self._parse_pdf
        if ext in self.WEB_EXTENSIONS:
            return self._parse_html
        if ext in {".odt", ".ods", ".odp"}:
            return self._parse_odf
        if ext == ".rtf":
            return self._parse_rtf
        if ext in self.DATA_EXTENSIONS:
            return self._parse_data
        return None

    def supported_extensions(self) -> List[str]:
        exts = []
        exts.extend(self.TEXT_EXTENSIONS)
        exts.extend(self.OFFICE_EXTENSIONS)
        exts.extend(self.WEB_EXTENSIONS)
        exts.extend(self.DATA_EXTENSIONS)
        return sorted(exts)

    # ──────────────────────────────────────────────
    #  PARSERS
    # ──────────────────────────────────────────────

    def _parse_text(self, path: Path, content: DocumentContent):
        raw = path.read_text("utf-8", errors="replace")
        content.raw_text = raw
        lines = raw.split("\n")
        cur_section = {"heading": "", "level": 1, "paragraphs": [], "tables": [], "lists": []}
        in_table = False
        table_buffer = []
        prev_line = ""
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if in_table and table_buffer:
                    t = self._parse_md_table(table_buffer)
                    if t:
                        content.tables.append(t)
                    cur_section["tables"].append(content.tables[-1] if content.tables else {})
                    table_buffer = []
                in_table = False
                continue
            # Heading: ATX or Setext
            heading_match = re.match(r'^(#{1,6})\s+(.+)', stripped)
            is_setext = False
            if not heading_match:
                setext_m = re.match(r'^(.+)\n([-=]+)$', prev_line + "\n" + stripped) if prev_line else None
                if setext_m:
                    heading_match = setext_m
                    is_setext = True
            if heading_match:
                if cur_section["paragraphs"] or cur_section["tables"] or cur_section["lists"]:
                    content.sections.append(cur_section)
                if is_setext:
                    dedup = heading_match.group(1).strip()
                    # Remove duplicate paragraph that was already added line N
                    if cur_section["paragraphs"] and cur_section["paragraphs"][-1] == dedup:
                        cur_section["paragraphs"].pop()
                    level = 1 if heading_match.group(2)[0] == "=" else 2
                    head_text = dedup
                else:
                    level = len(heading_match.group(1))
                    head_text = heading_match.group(2).strip()
                cur_section = {"heading": head_text, "level": level,
                               "paragraphs": [], "tables": [], "lists": []}
                prev_line = stripped
                continue
            # Table detection
            if "|" in stripped and "---" in stripped and prev_line and "|" in prev_line:
                table_buffer = [prev_line, stripped]
                in_table = True
                prev_line = stripped
                continue
            if in_table and "|" in stripped:
                table_buffer.append(stripped)
                prev_line = stripped
                continue
            # End of table
            if in_table and table_buffer:
                t = self._parse_md_table(table_buffer)
                if t:
                    content.tables.append(t)
                    cur_section["tables"].append(content.tables[-1])
                table_buffer = []
                in_table = False
            # List items
            if re.match(r'^[\s]*[-*+]\s', stripped) or re.match(r'^\s*\d+[.)]\s', stripped):
                item_text = re.sub(r'^[\s]*[-*+]\s|\s*\d+[.)]\s', '', stripped, count=1).strip()
                cur_section.setdefault("lists", []).append([item_text])
                prev_line = stripped
                continue
            # Regular paragraph
            cur_section["paragraphs"].append(stripped)
            prev_line = stripped
        # Flush last table
        if in_table and table_buffer:
            t = self._parse_md_table(table_buffer)
            if t:
                content.tables.append(t)
                cur_section["tables"].append(content.tables[-1])
        if cur_section["paragraphs"] or cur_section["tables"] or cur_section["lists"]:
            content.sections.append(cur_section)
        if not content.sections:
            content.paragraphs = [l.strip() for l in lines if l.strip()]

    @staticmethod
    def _parse_md_table(lines: list) -> dict:
        """Parse markdown table lines into {headers, rows}."""
        if len(lines) < 2:
            return {}
        headers = [c.strip() for c in lines[0].split("|") if c.strip()]
        rows = []
        for line in lines[2:]:
            cells = [c.strip() for c in line.split("|")]
            # Remove first/last if empty (leading/trailing |)
            if cells and cells[0] == "":
                cells = cells[1:]
            if cells and cells[-1] == "":
                cells = cells[:-1]
            if cells and any(c.strip() for c in cells):
                rows.append(cells[:len(headers)] if len(headers) else cells)
        return {"headers": headers, "rows": rows}

    def _parse_docx(self, path: Path, content: DocumentContent):
        from docx import Document as DocxDocument
        doc = DocxDocument(path)
        content.metadata = {
            "author": doc.core_properties.author or "",
            "created": str(doc.core_properties.created or ""),
            "modified": str(doc.core_properties.modified or ""),
        }
        elements = []
        for para in doc.paragraphs:
            style = para.style.name.lower() if para.style else ""
            text = para.text.strip()
            if not text:
                elements.append(("spacer", ""))
                continue
            if style.startswith("heading") or "titre" in style or "title" in style:
                level = 1
                m = re.search(r'heading\s*(\d+)', style)
                if m:
                    level = int(m.group(1))
                elements.append(("heading", text, min(level, 6)))
            else:
                elements.append(("para", text))
        # Tables
        tables_text = []
        for table in doc.tables:
            headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
            rows = []
            for row in table.rows[1:]:
                rows.append([cell.text.strip() for cell in row.cells])
            content.tables.append({"headers": headers, "rows": rows})
            tables_text.append(DocumentContent._table_to_text({"headers": headers, "rows": rows}))
        # Reconstruire sections
        cur = {"heading": "", "level": 1, "paragraphs": [], "tables": [], "lists": []}
        for el in elements:
            if el[0] == "spacer":
                if cur["paragraphs"] or cur["tables"]:
                    content.sections.append(cur)
                    cur = {"heading": "", "level": 1, "paragraphs": [], "tables": [], "lists": []}
                continue
            if el[0] == "heading":
                if cur["paragraphs"] or cur["tables"]:
                    content.sections.append(cur)
                cur = {"heading": el[1], "level": el[2], "paragraphs": [], "tables": [], "lists": []}
            else:
                cur["paragraphs"].append(el[1])
        if cur["paragraphs"] or cur["tables"]:
            content.sections.append(cur)
        # Reconstituer raw_text
        lines = []
        for sec in content.sections:
            if sec["heading"]:
                lines.append(f"{'#' * sec['level']} {sec['heading']}")
            lines.extend(sec["paragraphs"])
            for t in content.tables:
                lines.append(DocumentContent._table_to_text(t))
        lines.extend(tables_text)
        content.raw_text = "\n".join(lines)

    def _parse_pptx(self, path: Path, content: DocumentContent):
        from pptx import Presentation
        prs = Presentation(path)
        content.metadata = {"slides": str(len(prs.slides))}
        cur = {"heading": "", "level": 1, "paragraphs": [], "tables": [], "lists": []}
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    headers = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else []
                    rows = []
                    for row in table.rows[1:]:
                        rows.append([cell.text.strip() for cell in row.cells])
                    content.tables.append({"headers": headers, "rows": rows})
                    slide_texts.append(DocumentContent._table_to_text(content.tables[-1]))
            if slide_texts:
                cur["paragraphs"].append(f"[Slide {i+1}] " + " | ".join(slide_texts))
        if cur["paragraphs"] or cur["tables"]:
            content.sections.append(cur)
        content.raw_text = "\n".join(cur["paragraphs"])

    def _parse_xlsx(self, path: Path, content: DocumentContent):
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        sheets_text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in row_vals):
                    rows_data.append(row_vals)
            if not rows_data:
                continue
            headers = rows_data[0]
            data_rows = rows_data[1:]
            content.tables.append({"headers": headers, "rows": data_rows})
            sheets_text.append(f"## Sheet: {sheet}")
            sheets_text.append(DocumentContent._table_to_text(content.tables[-1]))
        content.raw_text = "\n".join(sheets_text)

    def _parse_pdf(self, path: Path, content: DocumentContent):
        from pypdf import PdfReader
        reader = PdfReader(path)
        content.page_count = len(reader.pages)
        content.metadata = {
            "pages": str(content.page_count),
            "author": reader.metadata.get("/Author", "") if reader.metadata else "",
        }
        cur = {"heading": "", "level": 1, "paragraphs": [], "tables": [], "lists": []}
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            lines = text.split("\n")
            for line in lines:
                stripped = line.strip()
                if not stripped:
                    continue
                # Heuristic: short uppercase or numbered lines may be headings
                if (stripped.isupper() and len(stripped) > 3 and len(stripped) < 80) or \
                   re.match(r'^(\d+\.\s+|[A-Z]\.\s+)', stripped):
                    if cur["paragraphs"] or cur["tables"]:
                        content.sections.append(cur)
                    cur = {"heading": stripped, "level": 2, "paragraphs": [], "tables": [], "lists": []}
                else:
                    cur["paragraphs"].append(stripped)
        if cur["paragraphs"] or cur["tables"]:
            content.sections.append(cur)
        # Reconstituer raw_text
        lines = []
        for sec in content.sections:
            if sec["heading"]:
                lines.append(sec["heading"])
            lines.extend(sec["paragraphs"])
        content.raw_text = "\n".join(lines)

    def _parse_html(self, path: Path, content: DocumentContent):
        raw = path.read_text("utf-8", errors="replace")
        # Strip tags via regex (lightweight, no bs4 dependency)
        text = re.sub(r'<style[^>]*>.*?</style>', '', raw, flags=re.DOTALL | re.I)
        text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.I)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        content.raw_text = text
        content.paragraphs = [p.strip() for p in text.split("\n") if p.strip()]

    def _parse_odf(self, path: Path, content: DocumentContent):
        from odf import text, teletype, office
        from odf.opendocument import load
        doc = load(path)
        all_paras = doc.getElementsByType(text.P)
        paragraphs = []
        for p in all_paras:
            t = teletype.extractText(p)
            if t.strip():
                paragraphs.append(t.strip())
        content.raw_text = "\n".join(paragraphs)
        content.paragraphs = paragraphs

    def _parse_rtf(self, path: Path, content: DocumentContent):
        raw = path.read_text("utf-8", errors="replace")
        text = re.sub(r'\\([a-z]+)([-0-9]+)?', ' ', raw)
        text = re.sub(r'\{|\}', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
        content.raw_text = text
        content.paragraphs = [p.strip() for p in text.split("\\par") if p.strip()]

    def _parse_data(self, path: Path, content: DocumentContent):
        ext = path.suffix.lower()
        raw = path.read_text("utf-8", errors="replace")
        if ext == ".csv":
            reader = csv.reader(io.StringIO(raw))
            rows = list(reader)
            if rows:
                content.tables.append({"headers": rows[0], "rows": rows[1:]})
                content.raw_text = DocumentContent._table_to_text(content.tables[-1])
        elif ext == ".json":
            try:
                data = json.loads(raw)
                content.raw_text = json.dumps(data, indent=2, ensure_ascii=False)
                content.paragraphs = [content.raw_text]
            except json.JSONDecodeError:
                content.raw_text = raw
                content.paragraphs = [raw]
        elif ext in {".xml", ".tsv"}:
            content.raw_text = raw
            if ext == ".tsv":
                reader = csv.reader(io.StringIO(raw), delimiter="\t")
                rows = list(reader)
                if rows:
                    content.tables.append({"headers": rows[0], "rows": rows[1:]})
            else:
                text = re.sub(r'<[^>]+>', ' ', raw)
                text = re.sub(r'\s+', ' ', text).strip()
                content.paragraphs = [text]
        else:
            content.raw_text = raw
            content.paragraphs = [raw]

    # ──────────────────────────────────────────────
    #  REGISTER DETECTION
    # ──────────────────────────────────────────────

    def _detect_register(self, text: str) -> str:
        """Détecte le registre dominant d'un document (scientific, legal, financial, etc.)."""
        if not text.strip():
            return "general"
        lower = text.lower()
        scores = Counter()
        for register, pattern in _REGISTER_KW_PATTERNS.items():
            count = len(pattern.findall(lower))
            if count:
                scores[register] += count * 2
        for register, pattern in _REGISTER_SENTINEL_PATTERNS.items():
            count = len(pattern.findall(text))
            if count:
                scores[register] += count * 3
        if not scores:
            return "general"
        best = scores.most_common(1)[0]
        return best[0] if best[1] >= 3 else "general"

    # ──────────────────────────────────────────────
    #  COMPRESSION
    # ──────────────────────────────────────────────

    def _compress_light(self, content: DocumentContent, category: str, mode: str = "light") -> str:
        """Structure-preserving compression. mode is passed to optimizer ('light', 'balanced', 'aggressive')."""
        parts = []
        sections = content.sections if content.sections else [
            {"heading": "", "level": 1, "paragraphs": content.paragraphs, "tables": content.tables, "lists": content.lists}
        ]
        for sec in sections:
            h = sec.get("heading", "")
            if h.lower().strip() in {"table of contents", "table des matières", "table des matieres",
                                      "sommaire", "index", "references", "références",
                                      "bibliography", "bibliographie", "copyright", "disclaimer"}:
                continue
            if h:
                parts.append(f"{'#' * min(sec.get('level',1),6)} {h}")
            for p in sec.get("paragraphs", []):
                if self._is_noise(p):
                    continue
                compressed = self._compress_paragraph(p, category, mode=mode)
                if compressed:
                    parts.append(compressed)
            for t in sec.get("tables", []):
                parts.append(self._compress_table(t, "light"))
            for lst in sec.get("lists", []):
                for item in lst:
                    ci = self._compress_paragraph(item, category, mode=mode)
                    parts.append(f"  - {ci}" if ci else f"  - {item}")
        return "\n".join(parts)

    def _compress_aggressive(self, content: DocumentContent, category: str) -> str:
        """Aggressive: ne garde que la substantifique moelle — titres, phrases-clés, données résumées."""
        parts = []
        noise_headings = {
            "table of contents", "table des matières", "table des matieres",
            "sommaire", "index", "references", "références", "bibliography",
            "bibliographie", "copyright", "disclaimer",
        }
        sections = content.sections if content.sections else [
            {"heading": "", "level": 1, "paragraphs": content.paragraphs, "tables": content.tables, "lists": content.lists}
        ]
        for sec in sections:
            h = sec.get("heading", "")
            if h.lower().strip() in noise_headings:
                continue
            if h:
                parts.append(h)
            for p in sec.get("paragraphs", []):
                if self._is_noise(p):
                    continue
                key_sentences = self._key_sentences(p, top_k=2)
                if key_sentences:
                    s = self._compress_paragraph(" ".join(key_sentences), category, mode="aggressive")
                    if s:
                        parts.append(s)
            for t in sec.get("tables", []):
                ct = self._compress_table(t, "aggressive")
                if ct:
                    parts.append(ct)
            for lst in sec.get("lists", []):
                for item in lst:
                    ci = self._compress_paragraph(item, category, mode="light")
                    parts.append(f"  - {ci}" if ci else "")
        cleaned = []
        for p in parts:
            p = re.sub(r'\s+', ' ', p).strip()
            if p:
                cleaned.append(p)
        return "\n".join(cleaned)

    def _compress_spc(self, content: DocumentContent, category: str, mode: str = "max") -> str:
        """Compression via SPC pipeline. mode: 'max' or 'industrial'."""
        try:
            from .spc.pipeline import SPC as SPCCompiler
            from .spc.profiles import MAX as SPC_MAX, INDUSTRIAL as SPC_INDUSTRIAL
            profile = SPC_INDUSTRIAL if mode == "industrial" else SPC_MAX
            compiler = SPCCompiler(profile=profile)
            result = compiler.compile(content.raw_text)
            return result.compressed
        except Exception:
            pass
        return content.raw_text

    @staticmethod
    def _is_noise(text: str) -> bool:
        """Détecte les lignes de bruit (copyright, page number, header/footer, TOC)."""
        t = text.strip().lower()
        if not t:
            return True
        # Page numbers seuls
        if re.match(r'^\d{1,4}$', t):
            return True
        if re.match(r'^page\s+\d+|p\.\s*\d+|-\s*\d+\s*-$', t, re.I):
            return True
        # Copyright, disclaimer, confidential
        if re.search(r'(copyright|©|all\s+rights?\s+reserved|tous\s+droits?\s+r[ée]serv[ée]s?'
                     r'|confidential|confidentiel|disclaimer|proprietary|propri[ée]taire'
                     r'|classified|classifi[ée])', t, re.I):
            return True
        # Très court (sauf si chiffres significatifs)
        word_count = len(t.split())
        if word_count < 3 and not re.search(r'\d+', t):
            return True
        return False

    def _compress_paragraph(self, text: str, category: str, mode: str = "light") -> str:
        """Compresse un paragraphe via l'optimiseur de prompts.
        mode: 'light', 'balanced', 'aggressive', 'max', or 'industrial'.
        """
        if len(text.split()) < 5:
            return text
        result = self.optimizer.optimize(text, category=category)
        return result.get(mode, {}).get("prompt", text).strip()

    def _compress_table(self, table: Dict, mode: str) -> str:
        headers = table.get("headers", [])
        rows = table.get("rows", [])
        if not rows and not headers:
            return ""
        if mode == "aggressive":
            # Ultra compact: headers → key values
            n = len(rows)
            if headers:
                line = " | ".join(headers)
                if n == 1:
                    line += " | " + " | ".join(str(c) for c in rows[0])
                else:
                    line += f" ({n} rows)"
                return line
            return f"Table: {n} rows, {len(headers)} cols" if headers else f"Table: {n} rows"
        # Light: markdown table compact
        lines = []
        if headers:
            lines.append("| " + " | ".join(headers) + " |")
        for row in rows[:20]:  # limite à 20 lignes
            lines.append("| " + " | ".join(str(c)[:60] for c in row) + " |")
        if len(rows) > 20:
            lines.append(f"| ... ({len(rows) - 20} more rows) |")
        return "\n".join(lines)

    def _key_sentences(self, text: str, top_k: int = 2) -> List[str]:
        """Extrait les phrases à plus haute densité informationnelle."""
        sentences = re.split(r'(?<=[.!?])\s+', text)
        if len(sentences) <= top_k + 1:
            return sentences
        n = len(sentences)
        scored = []
        for idx, s in enumerate(sentences):
            words = s.split()
            if len(words) < 3:
                continue
            long_words = sum(1 for w in words if len(w) > 6)
            numbers = len(re.findall(r'\d+', s))
            specific = len(re.findall(r'\b[A-Z][a-z]{2,}\b', s))
            score = long_words * 2 + numbers * 3 + specific
            if idx == 0:
                score += 2
            if idx == n - 1:
                score += 2
            scored.append((score, s))
        scored.sort(key=lambda x: -x[0])
        return [s for _, s in scored[:top_k]]

    # ──────────────────────────────────────────────
    #  UTILITIES
    # ──────────────────────────────────────────────

    def format_detection(self, filename: str) -> str:
        ext = Path(filename).suffix.lower()
        parser = self._get_parser(ext)
        if parser:
            return ext.lstrip(".")
        return "unknown"

    def supported_formats_summary(self) -> List[str]:
        return [f".{f}" for f in sorted(set(
            e.lstrip(".") for e in self.supported_extensions()
        ))]


# ──────────────────────────────────────────────
#  WRAPPER API
# ──────────────────────────────────────────────

def analyze_document(filepath: str, filename: str = "") -> dict:
    """Wrapper pour analyse rapide d'un document."""
    da = DocumentAnalyzer()
    content = da.analyze(filepath, filename)
    register = da._detect_register(content.raw_text)
    tokens = count_tokens(content.raw_text, "gpt-4o")
    return {
        "filename": filename or os.path.basename(filepath),
        "format": content.detected_format,
        "register": register,
        "tokens": tokens,
        "chars": len(content.raw_text),
        "words": len(content.raw_text.split()),
        "sections": len(content.sections) or 1,
        "tables": len(content.tables),
        "paragraphs": len(content.paragraphs) or sum(
            len(s.get("paragraphs", [])) for s in content.sections
        ) or 1,
        "preview": content.raw_text[:500],
    }


def compress_document(filepath: str, filename: str = "", mode: str = "light",
                      category: str = None) -> dict:
    """Wrapper pour compression complète d'un document."""
    da = DocumentAnalyzer()
    content = da.analyze(filepath, filename)
    compressed_text, meta = da.compress(content, mode=mode, category=category)
    return {
        **meta,
        "filename": filename or os.path.basename(filepath),
        "format": content.detected_format,
        "original_text": content.raw_text,
        "compressed_text": compressed_text,
        "preview_original": content.raw_text[:500],
        "preview_compressed": compressed_text[:500],
    }


if __name__ == "__main__":
    import sys
    da = DocumentAnalyzer()
    print(f"Formats supportés: {', '.join(da.supported_formats_summary())}")
    print(f"Registres: {', '.join(REGISTER_KEYWORDS.keys())}")
